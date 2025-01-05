import os
import hashlib
from typing import Tuple
import re, mimetypes
from docx import Document
import pandas as pd
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader
import logging

# Configure logging for structured logging
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
ch = logging.StreamHandler()
ch.setLevel(logging.INFO)
formatter = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')
ch.setFormatter(formatter)
logger.addHandler(ch)

def hash_document(content: str) -> str:
    """Generate a hash for the given document content."""
    return hashlib.sha256(content.encode('utf-8')).hexdigest()

def clean_text(text):
    """Clean text by removing extra spaces, fixing punctuation, and standardizing newlines."""
    text = re.sub(r'\s+', ' ', text)  # Replace multiple spaces with a single space
    text = re.sub(r'(?<=[a-zA-Z])(?=[A-Z])', ' ', text)  # Add space between words when stuck together
    text = re.sub(r'(?<=[.,?!;])(?=[a-zA-Z])', ' ', text)  # Add space after punctuation if missing
    text = re.sub(r'\.\s+', '.\n', text)  # Add newlines after sentences
    text = re.sub(r'(?<=:)\s+', '\n', text)  # Add newlines after colons
    text = text.strip()  # Remove leading/trailing spaces
    return text

def remove_null_bytes(text):
    """Remove null bytes from the text."""
    return text.replace('\x00', '')  # Removes null characters

def ensure_utf8_encoding(text):
    """Ensure the text is valid UTF-8 encoding."""
    try:
        return text.encode('utf-8').decode('utf-8')
    except UnicodeDecodeError:
        logger.error(f"Text contains invalid UTF-8 characters: {text[:30]}")
        return ""

def load_pdf(file_path):
    """Load PDF and return text from each page."""
    try:
        loader = PyPDFLoader(file_path)
        all_texts = []
        pages = loader.load_and_split()
        for page in pages:
            all_texts.append(page.page_content)
        return all_texts
    except Exception as e:
        logger.error(f"Error loading PDF {file_path}: {e}")
        return []

def load_word(file_path, chunk_size=10):
    """Load Word document and return text in chunks."""
    try:
        doc = Document(file_path)
        all_texts = []
        paragraphs = [paragraph.text for paragraph in doc.paragraphs]
        for i in range(0, len(paragraphs), chunk_size):
            chunk = "\n".join(paragraphs[i:i+chunk_size])
            all_texts.append(chunk)
        return all_texts
    except Exception as e:
        logger.error(f"Error loading Word document {file_path}: {e}")
        return []

def load_text(file_path, chunk_size=1024):
    """Load plain text file and return text in chunks, handle different encodings."""
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            all_texts = []
            while chunk := file.read(chunk_size):
                all_texts.append(chunk)
        return all_texts
    except UnicodeDecodeError:
        logger.warning(f"UTF-8 decoding failed for {file_path}, trying latin-1")
        with open(file_path, "r", encoding="latin-1") as file:
            all_texts = []
            while chunk := file.read(chunk_size):
                all_texts.append(chunk)
        return all_texts
    except Exception as e:
        logger.error(f"Error loading text file {file_path}: {e}")
        return []

def load_excel(file_path, chunk_size=1000):
    """Load Excel file and return text in chunks."""
    try:
        excel_data = pd.read_excel(file_path, chunksize=chunk_size)
        all_texts = []
        for chunk in excel_data:
            all_texts.append(chunk.to_string(index=False))
        return all_texts
    except Exception as e:
        logger.error(f"Error loading Excel file {file_path}: {e}")
        return []

def load_csv(file_path, chunk_size=1000):
    """Load CSV file and return text in chunks."""
    try:
        csv_data = pd.read_csv(file_path, chunksize=chunk_size)
        all_texts = []
        for chunk in csv_data:
            all_texts.append(chunk.to_string(index=False))
        return all_texts
    except Exception as e:
        logger.error(f"Error loading CSV file {file_path}: {e}")
        return []

def load_pptx(file_path):
    """Load PowerPoint file and return text from each slide."""
    try:
        presentation = Presentation(file_path)
        all_texts = []
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)
            all_texts.append("\n".join(slide_text))
        return all_texts
    except Exception as e:
        logger.error(f"Error loading PowerPoint file {file_path}: {e}")
        return []

def load_hidden_documents(directory="hidden_docs"):
    """Load all supported files from a directory and return their content in chunks."""
    all_texts = []
    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        mime_type, _ = mimetypes.guess_type(file_path)

        # Handle different file types based on extension
        try:
            if filename.endswith((".pdf", ".PDF")):
                all_texts.extend(load_pdf(file_path))
            elif filename.endswith((".docx", ".DOCX")):
                all_texts.extend(load_word(file_path))
            elif filename.endswith((".txt", ".TXT")):
                all_texts.extend(load_text(file_path))
            elif filename.endswith(('.xlsx', '.xls')):
                all_texts.extend(load_excel(file_path))
            elif filename.endswith((".csv", ".CSV")):
                all_texts.extend(load_csv(file_path))
            elif filename.endswith((".pptx", ".PPTX", ".ppt", ".PPT")):
                all_texts.extend(load_pptx(file_path))
            else:
                logger.warning(f"Unsupported file type for file: {filename}")
        except Exception as e:
            logger.error(f"Error processing file {filename}: {e}")

    # Clean the collected texts
    cleaned_texts = [clean_text(text) for text in all_texts]

    # Remove null bytes and ensure UTF-8 encoding
    cleaned_texts = [remove_null_bytes(text) for text in cleaned_texts]
    #cleaned_texts = [ensure_utf8_encoding(text) for text in cleaned_texts]

    # Filter out non-printable characters at the text level, not the character level
    cleaned_texts = [text for text in cleaned_texts if text.isprintable()]

    logger.info(f"Loaded and cleaned {len(cleaned_texts)} text chunks from documents.")
    return cleaned_texts

def store_documents(supabase, documents):
    """Store document hashes and contents in Supabase."""
    try:
        documents_to_insert = []
        for doc_text in documents:
            content_hash = hashlib.sha256(doc_text.encode()).hexdigest()

            # Check if document already exists using hash
            response = supabase.table("documents").select("id").eq("hash", content_hash).execute()

            if response is None:
                logger.error(f"Error checking existing document: {response.error}")
                continue

            if response.data:
                logger.info(f"Document with hash '{content_hash[:8]}...' already exists, skipping.")
                continue

            # Append both hash and content for new documents
            documents_to_insert.append({
                "hash": content_hash,
                "content": doc_text  # Store the actual document content
            })

        if documents_to_insert:  # Only insert if there are new documents
            response = supabase.table("documents").insert(documents_to_insert).execute()

            if response.error:
                logger.error(f"Error storing documents: {response.error}")
            else:
                logger.info(f"{len(documents_to_insert)} new documents stored successfully.")

    except Exception as e:
        logger.exception(f"Error interacting with Supabase: {e}")
